# -------------------------------------------------------------------------------- NOTEBOOK-CELL: CODE
# -*- coding: utf-8 -*-
"""
split_and_chunk.py  —  Step 2 of 3 (delta-aware)

Downloads, extracts, and chunks only the documents that load_data.py
identified as new/changed (the 'preprocessed_dataset' it produced is
already delta-filtered — this script does not re-filter).

The only functional additions vs. the original are:
  - All credentials read from environment variables instead of being
    hardcoded (see required env vars below).
  - 'version_hash' and 'delta_status' carried through onto every chunk row,
    so indexing.py can write them into Elasticsearch and use them on the
    next run's delta check.
  - A clean early-exit if the delta produced zero rows to process.

Required environment variables:
  PLM_TOKEN_URL, PLM_CLIENT_ID, PLM_CLIENT_SECRET, PLM_USER, PLM_PASSWORD,
  PLM_API_BASE, PLM_API_KEY,
  S3_ENDPOINT, S3_BUCKET, S3_ACCESS_KEY, S3_SECRET_KEY
"""
import dataiku
import sys
import pandas as pd, numpy as np
from dataiku import pandasutils as pdu
from pandarallel import pandarallel

import dataiku
import pandas as pd
import os
import json
import xlrd
import uuid
import base64
import logging
import shutil
import zipfile
import traceback
from datetime import datetime
from pathlib import Path
from typing import Optional
from uuid import uuid4
from openpyxl import load_workbook
from transformers import AutoTokenizer
from langchain_text_splitters import RecursiveCharacterTextSplitter
import boto3
import requests
import urllib3

urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

pandarallel.initialize(nb_workers=8)

# -------------------------------------------------------------------------------- NOTEBOOK-CELL: CODE
# ---------------------------------------------------------------------------
# Logging
# ---------------------------------------------------------------------------
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s  %(levelname)-8s  %(message)s",
    datefmt="%Y-%m-%d %H:%M:%S",
)
log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Configuration — all secrets come from environment variables / Dataiku
# project variables. Set these under Project > Variables (or your platform's
# secret store), never hardcode credentials in source.
# ---------------------------------------------------------------------------
PLM_AUTH = {
    "token_url": os.environ["PLM_TOKEN_URL"],
    "client_id": os.environ["PLM_CLIENT_ID"],
    "client_secret": os.environ["PLM_CLIENT_SECRET"],
    "username": os.environ["PLM_USER"],
    "password": os.environ["PLM_PASSWORD"],
    "scope": "profile",
}

PLM_API = {
    "base_url": os.environ["PLM_API_BASE"],
    "api_key": os.environ["PLM_API_KEY"],
}

S3_CONFIG = {
    "endpoint": os.environ["S3_ENDPOINT"],
    "bucket": os.environ["S3_BUCKET"],
    "access_key": os.environ["S3_ACCESS_KEY"],
    "secret_key": os.environ["S3_SECRET_KEY"],
}

download_folder = dataiku.Folder("exaZd6aI")
image_folder    = dataiku.Folder("2mIFClXn")

PIPELINE_CONFIG = {
    "download_dir": Path(download_folder.get_path()),
    "image_dir":    Path(image_folder.get_path()),
    "chunk_size": 500,
    "chunk_overlap": 50,
    "supported_formats": {".docx", ".doc", ".pdf", ".pptx", ".ppt", ".xlsx", ".xls", ".xlsm"},
    "download_folder": download_folder,
    "image_folder":    image_folder,
}
# PIPELINE_CONFIG["download_dir"].mkdir(parents=True, exist_ok=True)
# PIPELINE_CONFIG["image_dir"].mkdir(parents=True, exist_ok=True)
# ---------------------------------------------------------------------------
# S3 client
# ---------------------------------------------------------------------------
s3 = boto3.client(
    "s3",
    endpoint_url=S3_CONFIG["endpoint"],
    aws_access_key_id=S3_CONFIG["access_key"],
    aws_secret_access_key=S3_CONFIG["secret_key"],
    verify=False,
)
log.info("S3 client initialised.")

# ---------------------------------------------------------------------------
# Token cache
# ---------------------------------------------------------------------------
_token_cache: dict = {}

# ===========================================================================
# Helpers
# ===========================================================================

# -------------------------------------------------------------------------------- NOTEBOOK-CELL: CODE
model_folder = dataiku.Folder("U2J79e39")
path = model_folder.get_path()

tokenizer = AutoTokenizer.from_pretrained(pretrained_model_name_or_path=path)


def custom_tokenizer_length(text):
    tokens = tokenizer.tokenize(text)
    return len(tokens)


text_splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=10,
                                               length_function=custom_tokenizer_length)


def get_plm_access_token(force_refresh: bool = False) -> str:
    if not force_refresh and _token_cache.get("access_token"):
        return _token_cache["access_token"]
    raw = f"{PLM_AUTH['client_id']}:{PLM_AUTH['client_secret']}"
    b64 = base64.b64encode(raw.encode()).decode()
    headers = {
        "Authorization": f"Basic {b64}",
        "Content-Type": "application/x-www-form-urlencoded",
    }
    payload = {
        "grant_type": "password",
        "username": PLM_AUTH["username"],
        "password": PLM_AUTH["password"],
        "scope": PLM_AUTH["scope"],
    }
    resp = requests.post(PLM_AUTH["token_url"], headers=headers, data=payload,
                         timeout=30, verify=False)
    resp.raise_for_status()
    token = resp.json()["access_token"]
    _token_cache["access_token"] = token
    log.info("PLM access token acquired.")
    return token


def plm_headers() -> dict:
    try:
        token = get_plm_access_token()
        return {
            "Content-Type": "application/json",
            "X-Gravitee-Api-Key": PLM_API["api_key"],
            "Authorization": f"Bearer {token}",
            "accept": "*/*",
        }
    except Exception as exc:
        log.error("Failed to build PLM headers: %s", exc)
        return {"Content-Type": "application/json", "X-Gravitee-Api-Key": PLM_API["api_key"]}


def get_primary_content_meta(wt_document_id: str) -> Optional[dict]:
    try:
        if not wt_document_id.startswith("OR:"):
            wt_document_id = f"OR:wt.doc.WTDocument:{wt_document_id}"
        url = f"{PLM_API['base_url']}/DocMgmt/Documents('{wt_document_id}')/PrimaryContent"
        resp = requests.get(url, headers=plm_headers(), timeout=30, verify=False)
        if resp.status_code == 401:
            log.warning("PLM token expired – refreshing …")
            get_plm_access_token(force_refresh=True)
            resp = requests.get(url, headers=plm_headers(), timeout=30, verify=False)
        resp.raise_for_status()
        return resp.json()
    except Exception as exc:
        log.error("Error fetching PrimaryContent for %s: %s", wt_document_id, exc)
        return None


def download_document(download_url: str, dest_path: Path) -> bool:
    try:
        resp = requests.get(download_url, headers=plm_headers(),
                            stream=True, timeout=120, verify=False)
        resp.raise_for_status()
        with dest_path.open("wb") as fh:
            for chunk in resp.iter_content(chunk_size=65536):
                fh.write(chunk)
        return True
    except Exception as exc:
        log.error("Download failed for %s: %s", download_url, exc)
        return False


def extract_text_docx(file_path: Path) -> tuple:
    from docx import Document as DocxDocument
    import xml.etree.ElementTree as ET
    import subprocess
    image_list = []
    image_output_dir = PIPELINE_CONFIG["image_dir"]
    try:
        # ── Detect true file type regardless of extension ──────────────────
        # Some .doc files are actually OOXML (.docx) or other formats
        # Try opening as docx first; if it fails, attempt format conversion
        actual_path = file_path
        converted_path = None

        try:
            doc = DocxDocument(str(file_path))
        except Exception as fmt_exc:
            log.warning(
                "File %s failed direct docx open (%s) — attempting libreoffice conversion.",
                file_path.name, fmt_exc
            )
            # ── Attempt conversion via LibreOffice ─────────────────────────
            try:
                convert_dir = file_path.parent
                result = subprocess.run(
                    [
                        "libreoffice", "--headless", "--convert-to", "docx",
                        "--outdir", str(convert_dir),
                        str(file_path),
                    ],
                    timeout=120,
                    capture_output=True,
                    text=True,
                )
                converted_path = convert_dir / f"{file_path.stem}.docx"
                if result.returncode != 0 or not converted_path.exists():
                    log.error(
                        "LibreOffice conversion failed for %s: %s",
                        file_path.name, result.stderr
                    )
                    return "", image_list
                actual_path = converted_path
                doc = DocxDocument(str(actual_path))
                log.info("Successfully converted %s → %s", file_path.name, actual_path.name)
            except FileNotFoundError:
                log.error(
                    "LibreOffice not found — cannot convert %s. "
                    "Install with: apt-get install libreoffice",
                    file_path.name
                )
                return "", image_list
            except subprocess.TimeoutExpired:
                log.error("LibreOffice conversion timed out for %s", file_path.name)
                return "", image_list
            except Exception as conv_exc:
                log.error("Conversion error for %s: %s", file_path.name, conv_exc)
                return "", image_list

        image_map = {}

        with zipfile.ZipFile(str(file_path), "r") as docx_zip:
            rels_path = "word/_rels/document.xml.rels"
            if rels_path in docx_zip.namelist():
                rels_xml = docx_zip.read(rels_path)
                rels_tree = ET.fromstring(rels_xml)
                ns = "http://schemas.openxmlformats.org/package/2006/relationships"
                image_type = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
                for rel in rels_tree.findall(f"{{{ns}}}Relationship"):
                    if rel.attrib.get("Type") == image_type:
                        r_id = rel.attrib["Id"]
                        target = rel.attrib["Target"]
                        full_target = f"word/{target}"
                        if full_target in docx_zip.namelist():
                            suffix = Path(target).suffix
                            unique_name = f"{uuid4().hex}{suffix}"
                            save_path = image_output_dir / unique_name
                            with docx_zip.open(full_target) as img_file:
                                with open(save_path, "wb") as out_file:
                                    shutil.copyfileobj(img_file, out_file)
                            image_map[r_id] = unique_name
                            image_list.append(unique_name)

        parts = []
        for para in doc.paragraphs:
            para_parts = []
            for run in para.runs:
                if run.text:
                    para_parts.append(run.text)
                blips = run._r.findall(
                    ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
                )
                for blip in blips:
                    r_embed = blip.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                    )
                    if r_embed and r_embed in image_map:
                        para_parts.append(f"[IMAGE:{image_map[r_embed]}]")
            if para_parts:
                parts.append("".join(para_parts))

        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    parts.append(" | ".join(row_texts))

        if converted_path and converted_path.exists():
            try:
                converted_path.unlink()
                log.info("Deleted temporary converted file: %s", converted_path.name)
            except Exception as cleanup_exc:
                log.warning("Failed to delete converted file %s: %s", converted_path.name, cleanup_exc)

        return "\n".join(parts), image_list

    except Exception as exc:
        log.error("Error extracting docx %s: %s", file_path, exc)
        return "", image_list



def extract_text_pdf(file_path: Path) -> tuple:
    import fitz  # PyMuPDF

    image_list = []
    image_output_dir = PIPELINE_CONFIG["image_dir"]
    parts = []

    try:
        doc = fitz.open(str(file_path))
        for page in doc:
            parts.append(page.get_text())
            for img_index, img in enumerate(page.get_images(full=True)):
                base_image = doc.extract_image(img[0])
                unique_name = f"{uuid4().hex}.{base_image['ext']}"
                save_path = image_output_dir / unique_name
                save_path.write_bytes(base_image["image"])
                image_list.append(unique_name)
        doc.close()
        return "\n".join(parts), image_list
    except Exception as exc:
        log.error("Error extracting PDF %s: %s", file_path, exc)
        return "", image_list


def extract_text_pptx(file_path: Path) -> tuple:
    from pptx import Presentation
    image_list = []
    image_output_dir = PIPELINE_CONFIG["image_dir"]
    parts = []
    try:
        prs = Presentation(str(file_path))

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []

            for shape in slide.shapes:
                # ── Extract Text ──────────────────────────────
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                # ── Extract Images ────────────────────────────
                if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE = 13
                    try:
                        image = shape.image
                        ext = image.ext  # e.g., 'png', 'jpeg'
                        unique_name = f"{uuid4().hex}.{ext}"
                        save_path = image_output_dir / unique_name
                        save_path.write_bytes(image.blob)
                        image_list.append(unique_name)
                    except Exception as img_exc:
                        log.warning(
                            "Error extracting image from slide %d: %s",
                            slide_num, img_exc
                        )

            # ── Extract Images from placeholders/grouped shapes ──
            for shape in slide.placeholders:
                try:
                    if shape.placeholder_format and hasattr(shape, "image"):
                        image = shape.image
                        ext = image.ext
                        unique_name = f"{uuid4().hex}.{ext}"
                        save_path = image_output_dir / unique_name
                        save_path.write_bytes(image.blob)
                        image_list.append(unique_name)
                except Exception:
                    pass  # placeholder has no image, skip

            if slide_text:
                parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_text))

        return "\n\n".join(parts), image_list

    except Exception as exc:
        log.error("Error extracting PPTX %s: %s", file_path, exc)
        return "", image_list

def read_xlsx(filepath):
    """Read XLSX file using openpyxl"""
    text = []
    try:
        wb = load_workbook(filepath, read_only=True, data_only=True)
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_text = []
            for row in ws.iter_rows(values_only=True):
                row_text = " | ".join(
                    str(cell) for cell in row
                    if cell is not None and str(cell).strip()
                )
                if row_text.strip():
                    sheet_text.append(row_text)
            if sheet_text:
                text.append(f"[Sheet: {sheet_name}]\n" + "\n".join(sheet_text))
        wb.close()
    except Exception as e:
        print(e)
    return "\n\n".join(text), []

def read_xls(filepath):
    """Read old XLS file using xlrd"""
    text = []
    try:
        wb = xlrd.open_workbook(filepath)
        for sheet_name in wb.sheet_names():
            ws = wb.sheet_by_name(sheet_name)
            sheet_text = []
            for row_num in range(ws.nrows):
                row = ws.row_values(row_num)
                row_text = " | ".join(
                    str(cell).strip() for cell in row
                    if str(cell).strip()
                )
                if row_text.strip():
                    sheet_text.append(row_text)
            if sheet_text:
                text.append(f"[Sheet: {sheet_name}]\n" + "\n".join(sheet_text))
    except Exception as e:
        print(f"Error reading XLS: {e}")
    return "\n\n".join(text), []

def extract_text(file_path: Path) -> tuple:
    ext = file_path.suffix.lower()
    if ext in (".docx", ".doc"):
        return extract_text_docx(file_path)
    elif ext == ".pdf":
        return extract_text_pdf(file_path)
    elif ext in ('.pptx', '.ppt'):
        return extract_text_pptx(file_path)
    elif ext == '.xlsx' or ext == '.xlsm':
        return read_xlsx(file_path)
    elif ext == '.xls':
        return read_xls(file_path)

    return "", []


def upload_images_to_s3(doc_id: str, image_list: list) -> tuple:
    uploaded_images = {}
    failed_images = []
    image_output_dir = PIPELINE_CONFIG["image_dir"]

    magic_bytes_map = {
        b'\xff\xd8\xff': ('image/jpeg', '.jpg'),
        b'\x89PNG\r\n\x1a\n': ('image/png', '.png'),
        b'GIF87a': ('image/gif', '.gif'),
        b'GIF89a': ('image/gif', '.gif'),
        b'BM': ('image/bmp', '.bmp'),
        b'II\x2a\x00': ('image/tiff', '.tiff'),
        b'MM\x00\x2a': ('image/tiff', '.tiff'),
        b'RIFF': ('image/webp', '.webp'),
    }

    for idx, imagepath in enumerate(image_list):
        try:
            image_path = image_output_dir / imagepath
            final_bytes = image_path.read_bytes()
            if not final_bytes:
                raise ValueError("Empty image file.")

            ext = image_path.suffix.lower() or ".png"
            content_type = f"image/{ext.lstrip('.')}"
            for magic, (mime, extension) in magic_bytes_map.items():
                if final_bytes.startswith(magic):
                    content_type = mime
                    ext = extension
                    break

            image_id = f"{idx}_{image_path.stem}{ext}"
            s3_key = f"production_document/images_v0/{doc_id}/{image_id}"

            s3.put_object(
                Bucket=S3_CONFIG["bucket"],
                Key=s3_key,
                Body=final_bytes,
                ContentType=content_type,
            )
            s3_url = f"{S3_CONFIG['endpoint']}/{S3_CONFIG['bucket']}/{s3_key}"
            uploaded_images[imagepath] = s3_url

#             try:
#                 PIPELINE_CONFIG["image_folder"].delete_path(image_path.name)
#                 log.info("Deleted file from image folder  path=%s", image_path.name)
#             except Exception:
#                 log.warning(
#                     "Failed to delete file from image folder  path=%s\n%s",
#                     image_path.name,
#                     traceback.format_exc(),
#                 )

        except Exception as exc:
            log.error("Image upload failed (idx=%d, doc=%s): %s", idx, doc_id, exc)
            failed_images.append({"index": idx, "reason": str(exc)})

    return uploaded_images, failed_images


def chunk_text(text: str, chunk_size: int, overlap: int) -> list:
    return text_splitter.split_text(text)


# ===========================================================================
# Per-document processing
# ===========================================================================

def delete_pipeline_dirs(pipeline_config: dict):
    """Clear all files in managed folders using Dataiku API"""
    for folder_key, label in [("download_folder", "download"), ("image_folder", "image")]:
        folder = pipeline_config[folder_key]
        try:
            for item in folder.list_paths_in_partition():
                folder.delete_path(item)
            log.info("Cleared managed folder: %s", label)
        except Exception as e:
            log.error("Failed to clear managed folder %s: %s", label, e)

def process_document(doc_row, run_id: str = str(uuid.uuid4())) -> list:  # run_id captured via default arg

    doc_row = doc_row.to_dict() if hasattr(doc_row, 'to_dict') else doc_row
    doc_id_raw = doc_row.get("id", "")
    doc_id = str(int(float(str(doc_id_raw).strip())))
    doc_num = doc_row.get("doc_number", "")
    filename = doc_row.get("filename", "") or ""
    # Carried through from load_data.py's delta check — needed by indexing.py
    # to (a) know whether to delete old chunks first and (b) stamp the new
    # chunks with the version that's now current in ES.
    version_hash = doc_row.get("version_hash", "")
    delta_status = doc_row.get("delta_status", "")
    # log.info("Processing  id=%-15s  number=%-15s  file=%s", doc_id, doc_num, filename)

    try:
        PIPELINE_CONFIG["download_dir"].mkdir(parents=True, exist_ok=True)
        PIPELINE_CONFIG["image_dir"].mkdir(parents=True, exist_ok=True)


        content_meta = get_primary_content_meta(doc_id)
        if not content_meta:
            log.warning("Skipping %s – no PrimaryContent metadata.", doc_id)
            return []

        file_name = content_meta.get("FileName", filename)
        download_url = (content_meta.get("Content") or {}).get("URL")
        if not download_url:
            log.warning("Skipping %s – no download URL.", doc_id)
            return []

        file_ext = Path(file_name).suffix.lower()
        if file_ext not in PIPELINE_CONFIG["supported_formats"]:
            log.info("Skipping %s – unsupported format '%s'.", doc_id, file_ext)
            return []

        # Download file
        local_path = PIPELINE_CONFIG["download_dir"] / f"{doc_id}{file_ext}"
        if not download_document(download_url, local_path):
            return []

        # Extract text + images
        full_text, image_list = extract_text(local_path)
        if not full_text.strip():
            log.warning("No text extracted from %s.", local_path)

        # Upload images to S3
        uploaded_map, _ = upload_images_to_s3(doc_id, image_list)

        try:
            PIPELINE_CONFIG["download_folder"].delete_path(local_path.name)
            log.info("Deleted file from managed folder  path=%s", local_path.name)
        except Exception:
            log.warning(
                "Failed to delete file from managed folder  path=%s\n%s",
                local_path.name,
                traceback.format_exc(),
            )

        # Chunk text and replace image placeholders with S3 URLs
        chunks = chunk_text(full_text, PIPELINE_CONFIG["chunk_size"], PIPELINE_CONFIG["chunk_overlap"])
        chunk_rows = []

        for chunk_seq, text in enumerate(chunks):
            for placeholder, s3_url in uploaded_map.items():
                text = text.replace(placeholder, s3_url)

            chunk_rows.append({
                "run_id": run_id,
                "doc_id": doc_id,
                "doc_number": doc_num,
                "filename": file_name,
                "filepath": f"https://plmpublishing.icp.infineon.com/api/download-pdf/{doc_num}",
                "filetype": file_ext.lstrip("."),
                "chunk_seq": chunk_seq,
                "document_type": doc_row.get("type", ""),
                "document_sub_type": doc_row.get("doc_sub_type", ""),
                "location": doc_row.get("infineon_loc", ""),
                "chunk_text": text,
                "images_count": len(uploaded_map),
                "chunked_at": datetime.utcnow().isoformat(),
                "version_hash": version_hash,
                "delta_status": delta_status,
            })

        # log.info("doc_id=%s → %d chunks, %d images uploaded.", doc_id, len(chunk_rows), len(uploaded_map))
        return chunk_rows

    except Exception as exc:
        log.error("Unhandled error processing %s: %s\n%s", doc_id, exc, traceback.format_exc())
        return []

# -------------------------------------------------------------------------------- NOTEBOOK-CELL: CODE
# Read recipe inputs.
# This dataset is already delta-filtered by load_data.py — it only contains
# documents that are new or changed vs. what's currently indexed in ES.
preprocessed_dataset = dataiku.Dataset(
    "preprocessed_dataset")
df = preprocessed_dataset.get_dataframe()

# NOTE: the original script had `df = df.iloc[10:]` here, which silently
# dropped the first 10 rows on every run (looked like leftover debug code).
# Removed — with delta filtering already narrowing the dataset upstream,
# skipping rows here would just lose documents. Re-add deliberately if there
# was a real reason for it (e.g. a known-bad header block in some other feed).

log.info("Documents to process this run (post-delta): %d", len(df))

if df.empty:
    log.info("No new or changed documents this run — nothing to download/chunk.")
    production_split_chunk_df = pd.DataFrame(columns=[
        "run_id", "doc_id", "doc_number", "filename", "filepath", "filetype",
        "chunk_seq", "document_type", "document_sub_type", "location",
        "chunk_text", "images_count", "chunked_at", "version_hash", "delta_status",
    ])
else:
    # -------------------------------------------------------------------------------- NOTEBOOK-CELL: CODE
    delete_pipeline_dirs(PIPELINE_CONFIG)

    # -------------------------------------------------------------------------------- NOTEBOOK-CELL: CODE
    results = df.parallel_apply(process_document, axis=1)

    all_chunks = [chunk for chunk_list in results if chunk_list for chunk in chunk_list]

    # -------------------------------------------------------------------------------- NOTEBOOK-CELL: CODE
    production_split_chunk_df = pd.DataFrame(all_chunks)

# Write recipe outputs
production_split_chunk = dataiku.Dataset("production_split_chunk")
production_split_chunk.write_with_schema(production_split_chunk_df)
